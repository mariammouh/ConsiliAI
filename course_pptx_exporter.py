"""
Professional Course PowerPoint Exporter — Robust Dynamic Layout Edition
Generates beautifully designed, adaptive PPTX files per lesson.
No text truncation. No overflow. Dynamic vertical positioning.
"""

import os
import re
import math
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# =============================================================================
# DESIGN TOKENS
# =============================================================================

@dataclass(frozen=True)
class DesignTokens:
    PRIMARY: str = "#1E40AF"
    ACCENT: str = "#3B82F6"
    BACKGROUND: str = "#F8FAFC"
    DARK_TEXT: str = "#1E293B"
    SECONDARY_GRAY: str = "#64748B"
    SUCCESS: str = "#10B981"
    WARNING: str = "#F59E0B"
    WHITE: str = "#FFFFFF"
    LIGHT_GRAY: str = "#E2E8F0"
    CARD_BG: str = "#FFFFFF"
    LIGHT_BLUE: str = "#DBEAFE"
    LIGHTER_BLUE: str = "#93C5FD"
    VERY_LIGHT_BLUE: str = "#E0F2FE"

    TITLE_FONT: str = "Calibri"
    BODY_FONT: str = "Calibri"

    SLIDE_WIDTH: Inches = Inches(13.333)
    SLIDE_HEIGHT: Inches = Inches(7.5)

    LEFT_MARGIN: Inches = Inches(0.6)
    RIGHT_MARGIN: Inches = Inches(0.6)
    TOP_SAFE: Inches = Inches(0.35)
    BOTTOM_SAFE: Inches = Inches(0.35)

    HEADER_HEIGHT: Inches = Inches(0.10)
    FOOTER_HEIGHT: Inches = Inches(0.30)

    CONTENT_TOP: Inches = Inches(0.35)
    CONTENT_BOTTOM: Inches = Inches(7.15)

    CARD_PADDING: Inches = Inches(0.22)
    GAP_SMALL: Inches = Inches(0.12)
    GAP_MEDIUM: Inches = Inches(0.20)
    GAP_LARGE: Inches = Inches(0.30)

    COVER_TITLE: Pt = Pt(50)
    COVER_SUBTITLE: Pt = Pt(24)
    DIVIDER_TITLE: Pt = Pt(44)
    SLIDE_TITLE: Pt = Pt(28)
    BODY_TEXT: Pt = Pt(16)
    BODY_SMALL: Pt = Pt(14)
    FOOTER_TEXT: Pt = Pt(10)
    BADGE_TEXT: Pt = Pt(12)

    MAX_BULLETS_PER_SLIDE: int = 4
    MAX_CHARS_PER_BULLET: int = 65
    MAX_QUESTIONS_PER_SLIDE: int = 3
    MAX_REFS_PER_SLIDE: int = 3
    MAX_TERMS_PER_SLIDE: int = 4


TOKENS = DesignTokens()


# =============================================================================
# UTILITIES
# =============================================================================

def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def to_inches(val) -> float:
    """Convert Inches/Length/int(EMU) to float inches."""
    if hasattr(val, 'inches'):
        return float(val.inches)
    if hasattr(val, 'pt'):
        return float(val.pt) / 72.0
    return float(val) / 914400.0


def set_fill(shape, color: str):
    r, g, b = hex_to_rgb(color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)


def no_outline(shape):
    shape.line.fill.background()


def add_rounded_rect(slide, left, top, width, height, fill=TOKENS.CARD_BG, outline=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(s, fill)
    if outline:
        s.line.color.rgb = RGBColor(*hex_to_rgb(outline))
        s.line.width = Pt(1)
    else:
        no_outline(s)
    return s


def add_accent_bar(slide, left, top, width, height, color=TOKENS.ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_fill(s, color)
    no_outline(s)
    return s


def add_bg(slide, color=TOKENS.BACKGROUND):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                TOKENS.SLIDE_WIDTH, TOKENS.SLIDE_HEIGHT)
    set_fill(s, color)
    no_outline(s)
    return s


def add_header(slide, color=TOKENS.PRIMARY, height=TOKENS.HEADER_HEIGHT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                TOKENS.SLIDE_WIDTH, height)
    set_fill(s, color)
    no_outline(s)
    return s


def add_footer(slide, num: int, course_title: str = ""):
    y_line = TOKENS.SLIDE_HEIGHT - TOKENS.FOOTER_HEIGHT - Inches(0.01)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y_line,
                                   TOKENS.SLIDE_WIDTH, Inches(0.01))
    set_fill(line, TOKENS.LIGHT_GRAY)
    no_outline(line)

    if course_title:
        box = slide.shapes.add_textbox(TOKENS.LEFT_MARGIN,
            TOKENS.SLIDE_HEIGHT - TOKENS.FOOTER_HEIGHT + Inches(0.05),
            Inches(6), Inches(0.2))
        p = box.text_frame.paragraphs[0]
        p.text = course_title
        p.font.size = TOKENS.FOOTER_TEXT
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.SECONDARY_GRAY))
        p.font.name = TOKENS.BODY_FONT

    box = slide.shapes.add_textbox(TOKENS.SLIDE_WIDTH - Inches(1.0),
        TOKENS.SLIDE_HEIGHT - TOKENS.FOOTER_HEIGHT + Inches(0.05),
        Inches(0.6), Inches(0.2))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(num)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = TOKENS.FOOTER_TEXT
    p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.SECONDARY_GRAY))
    p.font.name = TOKENS.BODY_FONT
    p.font.bold = True


# =============================================================================
# TEXT METRICS — CONSERVATIVE ESTIMATES
# =============================================================================

class TextMetrics:
    """Conservative text dimension estimates to prevent truncation."""

    @staticmethod
    def cpi(font_size_pt: float) -> float:
        return max(3.0, 95.0 / font_size_pt)

    @staticmethod
    def line_h(font_size_pt: float, spacing: float = 1.4) -> float:
        return (font_size_pt / 72.0) * spacing

    @staticmethod
    def lines(text: str, font_size_pt: float, width_inches: float) -> int:
        if not text or not text.strip():
            return 0
        cpl = max(1, int(width_inches * TextMetrics.cpi(font_size_pt)))
        words = text.split()
        lines = 0
        cur = 0
        for w in words:
            wl = len(w)
            if cur + wl + 1 > cpl:
                lines += 1
                cur = wl
            else:
                cur += wl + 1
        if cur > 0:
            lines += 1
        return lines

    @staticmethod
    def height(text: str, font_size_pt: float, width_inches: float, spacing: float = 1.4) -> float:
        n = TextMetrics.lines(text, font_size_pt, width_inches)
        if n == 0:
            return 0.0
        return n * TextMetrics.line_h(font_size_pt, spacing)

    @staticmethod
    def bullet_height(text: str, font_size_pt: float, width_inches: float) -> float:
        w = max(1.0, width_inches - 0.30)
        h = TextMetrics.height(text, font_size_pt, w, 1.4)
        return max(h, TextMetrics.line_h(font_size_pt, 1.4))

    @staticmethod
    def list_height(bullets: List[str], font_size_pt: float, width_inches: float, gap: float = 0.12) -> float:
        total = 0.0
        for b in bullets:
            total += TextMetrics.bullet_height(b, font_size_pt, width_inches) + gap
        return max(0.0, total - gap)


# =============================================================================
# CONTENT SPLITTING
# =============================================================================

class ContentSplitter:
    @staticmethod
    def to_bullets(text: str, max_bullets: int = TOKENS.MAX_BULLETS_PER_SLIDE,
                   max_chars: int = TOKENS.MAX_CHARS_PER_BULLET) -> List[List[str]]:
        if not text:
            return []

        raw = re.split(r'(?<=[.!?;])\s+', text)
        bullets = []
        current = ""
        for chunk in raw:
            chunk = chunk.strip()
            if not chunk:
                continue
            if len(current) + len(chunk) + 1 <= max_chars:
                current += " " + chunk if current else chunk
            else:
                if current:
                    bullets.append(current.strip())
                current = chunk
        if current:
            bullets.append(current.strip())

        refined = []
        for b in bullets:
            if len(b) <= max_chars:
                refined.append(b)
            else:
                words = b.split()
                cur = ""
                for w in words:
                    if len(cur) + len(w) + 1 <= max_chars:
                        cur += " " + w if cur else w
                    else:
                        if cur:
                            refined.append(cur.strip())
                        cur = w
                if cur:
                    refined.append(cur.strip())

        groups = []
        group = []
        for b in refined:
            group.append(b)
            if len(group) >= max_bullets:
                groups.append(group)
                group = []
        if group:
            groups.append(group)
        return groups if groups else [[]]

    @staticmethod
    def to_paragraphs(text: str, max_chars: int = 450) -> List[str]:
        if not text:
            return []
        paras = [p.strip() for p in text.split("\n") if p.strip()]
        if not paras:
            paras = [text.strip()]
        result = []
        for para in paras:
            if len(para) <= max_chars:
                result.append(para)
            else:
                sentences = re.split(r'(?<=[.!?])\s+', para)
                chunk = ""
                for s in sentences:
                    if len(chunk) + len(s) + 1 <= max_chars:
                        chunk += " " + s if chunk else s
                    else:
                        if chunk:
                            result.append(chunk.strip())
                        chunk = s
                if chunk:
                    result.append(chunk.strip())
        return result

    @staticmethod
    def chunk_dict(items: Dict[str, str], per_slide: int = TOKENS.MAX_TERMS_PER_SLIDE) -> List[Dict[str, str]]:
        if not items:
            return []
        pairs = list(items.items())
        return [dict(pairs[i:i+per_slide]) for i in range(0, len(pairs), per_slide)]

    @staticmethod
    def chunk_list(items: List[Any], per_slide: int = 4) -> List[List[Any]]:
        if not items:
            return []
        return [items[i:i+per_slide] for i in range(0, len(items), per_slide)]


# =============================================================================
# LAYOUT ENGINE
# =============================================================================

class Layout:
    def __init__(self, slide, course_title: str = "", slide_number: int = 1):
        self.slide = slide
        self.course_title = course_title
        self.slide_number = slide_number
        self.y = to_inches(TOKENS.CONTENT_TOP)
        self.left = TOKENS.LEFT_MARGIN
        self.right = TOKENS.RIGHT_MARGIN
        self.width_inches = to_inches(TOKENS.SLIDE_WIDTH) - to_inches(TOKENS.LEFT_MARGIN) - to_inches(TOKENS.RIGHT_MARGIN)
        self.bottom = to_inches(TOKENS.CONTENT_BOTTOM)

    def remaining(self) -> float:
        return self.bottom - self.y

    def advance(self, h: float):
        self.y += h


# =============================================================================
# RENDERERS — ANTI-TRUNCATION DESIGN
# =============================================================================

def render_title(layout: Layout, text: str, font_size: Pt = TOKENS.SLIDE_TITLE,
                 color: str = TOKENS.DARK_TEXT, bold: bool = True):
    h = TextMetrics.height(text, font_size.pt, layout.width_inches, 1.25)
    h = max(h, font_size.pt / 72.0 * 1.25)
    box = layout.slide.shapes.add_textbox(
        layout.left, Inches(layout.y), Inches(layout.width_inches), Inches(h + 0.15)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = RGBColor(*hex_to_rgb(color))
    p.font.name = TOKENS.TITLE_FONT
    p.font.bold = bold
    layout.advance(h + 0.15)


def render_accent_bar(layout: Layout, width_inches: float = 1.0, color: str = TOKENS.ACCENT):
    add_accent_bar(layout.slide, layout.left, Inches(layout.y), Inches(width_inches), Inches(0.04), color)
    layout.advance(0.04 + to_inches(TOKENS.GAP_SMALL))


def render_bullet_card(layout: Layout, bullets: List[str],
                       fill: str = TOKENS.CARD_BG, outline: str = TOKENS.LIGHT_GRAY,
                       font_size: Pt = TOKENS.BODY_TEXT, bullet_color: str = TOKENS.ACCENT):
    if not bullets:
        return

    pad = to_inches(TOKENS.CARD_PADDING)
    text_w = layout.width_inches - 2 * pad

    text_h = TextMetrics.list_height(bullets, font_size.pt, text_w, 0.12)
    card_h = max(2.0, text_h + 2 * pad + 0.30)

    max_card_h = layout.remaining() - to_inches(TOKENS.GAP_MEDIUM)
    card_h = min(card_h, max_card_h)

    add_rounded_rect(layout.slide, layout.left, Inches(layout.y),
                     Inches(layout.width_inches), Inches(card_h), fill, outline)

    text_box_h = min(Inches(layout.remaining()), Inches(5.0))
    text_box = layout.slide.shapes.add_textbox(
        layout.left + TOKENS.CARD_PADDING, Inches(layout.y + pad),
        Inches(text_w), text_box_h
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = font_size
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
        p.font.name = TOKENS.BODY_FONT
        p.space_after = Pt(10)
        p.line_spacing = 1.35
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor(*hex_to_rgb(bullet_color))

    layout.advance(card_h + to_inches(TOKENS.GAP_MEDIUM))


def render_paragraph_card(layout: Layout, paragraphs: List[str],
                          fill: str = TOKENS.CARD_BG, outline: str = TOKENS.LIGHT_GRAY,
                          font_size: Pt = TOKENS.BODY_TEXT, italic: bool = False):
    if not paragraphs:
        return

    pad = to_inches(TOKENS.CARD_PADDING)
    text_w = layout.width_inches - 2 * pad

    text_h = 0.0
    for para in paragraphs:
        text_h += TextMetrics.height(para, font_size.pt, text_w, 1.35) + 0.12
    if paragraphs:
        text_h -= 0.12

    card_h = max(1.8, text_h + 2 * pad + 0.30)
    max_card_h = layout.remaining() - to_inches(TOKENS.GAP_MEDIUM)
    card_h = min(card_h, max_card_h)

    add_rounded_rect(layout.slide, layout.left, Inches(layout.y),
                     Inches(layout.width_inches), Inches(card_h), fill, outline)

    text_box_h = min(Inches(layout.remaining()), Inches(5.0))
    text_box = layout.slide.shapes.add_textbox(
        layout.left + TOKENS.CARD_PADDING, Inches(layout.y + pad),
        Inches(text_w), text_box_h
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para
        p.font.size = font_size
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
        p.font.name = TOKENS.BODY_FONT
        p.font.italic = italic
        p.space_after = Pt(10)
        p.line_spacing = 1.35

    layout.advance(card_h + to_inches(TOKENS.GAP_MEDIUM))


def render_numbered_card(layout: Layout, items: List[str],
                         fill: str = TOKENS.CARD_BG, outline: str = TOKENS.LIGHT_GRAY,
                         font_size: Pt = TOKENS.BODY_TEXT, number_color: str = TOKENS.ACCENT):
    if not items:
        return

    pad = to_inches(TOKENS.CARD_PADDING)
    text_w = layout.width_inches - 2 * pad

    text_h = 0.0
    for item in items:
        text_h += TextMetrics.bullet_height(item, font_size.pt, text_w) + 0.12
    if items:
        text_h -= 0.12

    card_h = max(2.0, text_h + 2 * pad + 0.30)
    max_card_h = layout.remaining() - to_inches(TOKENS.GAP_MEDIUM)
    card_h = min(card_h, max_card_h)

    add_rounded_rect(layout.slide, layout.left, Inches(layout.y),
                     Inches(layout.width_inches), Inches(card_h), fill, outline)

    text_box_h = min(Inches(layout.remaining()), Inches(5.0))
    text_box = layout.slide.shapes.add_textbox(
        layout.left + TOKENS.CARD_PADDING, Inches(layout.y + pad),
        Inches(text_w), text_box_h
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = font_size
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
        p.font.name = TOKENS.BODY_FONT
        p.space_after = Pt(10)
        p.line_spacing = 1.35
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor(*hex_to_rgb(number_color))
            p.runs[0].font.bold = True

    layout.advance(card_h + to_inches(TOKENS.GAP_MEDIUM))


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def build_title_slide(prs, course_title, lesson_title, module_title="", slide_number=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, TOKENS.PRIMARY, Inches(0.12))

    add_rounded_rect(slide, Inches(8.5), Inches(2.5), Inches(4), Inches(4), TOKENS.LIGHT_BLUE)
    add_rounded_rect(slide, Inches(11), Inches(1.2), Inches(0.8), Inches(0.8), TOKENS.WARNING)

    if course_title:
        box = slide.shapes.add_textbox(TOKENS.LEFT_MARGIN, Inches(2.1), Inches(8), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = course_title.upper()
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.ACCENT))
        p.font.name = TOKENS.BODY_FONT
        p.font.bold = True

    h = TextMetrics.height(lesson_title, TOKENS.COVER_TITLE.pt, to_inches(Inches(9)), 1.2)
    h = max(h, TOKENS.COVER_TITLE.pt / 72.0 * 1.2)
    box = slide.shapes.add_textbox(TOKENS.LEFT_MARGIN, Inches(2.6), Inches(9), Inches(h + 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = lesson_title
    p.font.size = TOKENS.COVER_TITLE
    p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
    p.font.name = TOKENS.TITLE_FONT
    p.font.bold = True

    bar_y = Inches(2.6 + h + 0.15)
    add_accent_bar(slide, TOKENS.LEFT_MARGIN, bar_y, Inches(1.4), Inches(0.05), TOKENS.ACCENT)

    if module_title:
        box = slide.shapes.add_textbox(TOKENS.LEFT_MARGIN, bar_y + Inches(0.15), Inches(8), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = module_title
        p.font.size = TOKENS.COVER_SUBTITLE
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.SECONDARY_GRAY))
        p.font.name = TOKENS.BODY_FONT

    add_footer(slide, slide_number, course_title)
    return slide


def build_objectives_slide(prs, objectives, slide_number, course_title=""):
    if not objectives:
        return None
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide)

    layout = Layout(slide, course_title, slide_number)
    render_title(layout, "Learning Objectives")
    render_accent_bar(layout, 1.2, TOKENS.ACCENT)

    card_width = Inches(5.6)
    gap_x = Inches(0.35)
    cols = 2
    pad = to_inches(TOKENS.CARD_PADDING)

    for i, obj in enumerate(objectives):
        row = i // cols
        col = i % cols
        x = layout.left + col * (card_width + gap_x)

        text_w = to_inches(card_width) - 2 * pad - 0.6
        text_h = TextMetrics.height(obj, TOKENS.BODY_TEXT.pt, text_w, 1.35)
        card_h = max(1.1, text_h + 2 * pad + 0.15)
        y = layout.y + row * (card_h + 0.2)

        add_rounded_rect(slide, x, Inches(y), card_width, Inches(card_h), TOKENS.CARD_BG, TOKENS.LIGHT_GRAY)

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), Inches(y + 0.25), Inches(0.42), Inches(0.42))
        set_fill(badge, TOKENS.ACCENT)
        no_outline(badge)

        num = slide.shapes.add_textbox(x + Inches(0.12), Inches(y + 0.28), Inches(0.42), Inches(0.38))
        p = num.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WHITE))
        p.font.name = TOKENS.BODY_FONT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        txt = slide.shapes.add_textbox(x + Inches(0.65), Inches(y + 0.15), card_width - Inches(0.8), Inches(card_h - 0.3))
        tf = txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = obj
        p.font.size = TOKENS.BODY_TEXT
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
        p.font.name = TOKENS.BODY_FONT

    rows = math.ceil(len(objectives) / cols)
    if rows > 0:
        layout.y += rows * 1.3
    add_footer(slide, slide_number, course_title)
    return slide


def build_context_slide(prs, problem, solution, difficulty, slide_number, course_title=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide)

    layout = Layout(slide, course_title, slide_number)
    render_title(layout, "Module Context")
    render_accent_bar(layout, 1.2, TOKENS.ACCENT)

    cards = [("Problem", problem, TOKENS.WARNING), ("Solution", solution, TOKENS.SUCCESS),
             ("Difficulty", difficulty, TOKENS.ACCENT)]
    cards = [(l, c, col) for l, c, col in cards if c]
    if not cards:
        add_footer(slide, slide_number, course_title)
        return slide

    n = len(cards)
    gap = Inches(0.25)
    card_w = Inches((layout.width_inches - (n - 1) * to_inches(gap)) / n)
    start_y = layout.y
    max_h = 0
    pad = to_inches(TOKENS.CARD_PADDING)

    for i, (label, content, color) in enumerate(cards):
        x = layout.left + i * (card_w + gap)
        text_w = to_inches(card_w) - 2 * pad
        label_h = TextMetrics.line_h(12, 1.3)
        content_h = TextMetrics.height(content, TOKENS.BODY_TEXT.pt, text_w, 1.35)
        card_h = max(2.2, label_h + content_h + 2 * pad + 0.3)
        max_h = max(max_h, card_h)

        add_rounded_rect(slide, x, Inches(start_y), card_w, Inches(card_h), TOKENS.CARD_BG, TOKENS.LIGHT_GRAY)

        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(start_y), card_w, Inches(0.06))
        set_fill(strip, color)
        no_outline(strip)

        lbl = slide.shapes.add_textbox(x + TOKENS.CARD_PADDING, Inches(start_y + 0.15), card_w - Inches(0.4), Inches(0.35))
        p = lbl.text_frame.paragraphs[0]
        p.text = label.upper()
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(*hex_to_rgb(color))
        p.font.name = TOKENS.BODY_FONT
        p.font.bold = True

        txt = slide.shapes.add_textbox(x + TOKENS.CARD_PADDING, Inches(start_y + 0.5), card_w - Inches(0.4), Inches(card_h - 0.7))
        tf = txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = TOKENS.BODY_TEXT
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
        p.font.name = TOKENS.BODY_FONT
        p.line_spacing = 1.3

    layout.advance(max_h + to_inches(TOKENS.GAP_MEDIUM))
    add_footer(slide, slide_number, course_title)
    return slide


def build_divider_slide(prs, title, subtitle="", slide_number=1, course_title=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.2), TOKENS.SLIDE_HEIGHT)
    set_fill(left, TOKENS.PRIMARY)
    no_outline(left)

    right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(0), TOKENS.SLIDE_WIDTH - Inches(4.2), TOKENS.SLIDE_HEIGHT)
    set_fill(right, TOKENS.BACKGROUND)
    no_outline(right)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), Inches(2.9), Inches(1.4), Inches(1.4))
    set_fill(circle, TOKENS.ACCENT)
    no_outline(circle)

    icon = slide.shapes.add_textbox(Inches(1.6), Inches(3.2), Inches(1.0), Inches(0.8))
    p = icon.text_frame.paragraphs[0]
    p.text = "§"
    p.font.size = Pt(44)
    p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WHITE))
    p.font.name = TOKENS.TITLE_FONT
    p.alignment = PP_ALIGN.CENTER

    h = TextMetrics.height(title, TOKENS.DIVIDER_TITLE.pt, to_inches(Inches(8.5)), 1.2)
    h = max(h, TOKENS.DIVIDER_TITLE.pt / 72.0 * 1.2)
    box = slide.shapes.add_textbox(Inches(4.8), Inches(2.7), Inches(8.5), Inches(h + 0.15))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = TOKENS.DIVIDER_TITLE
    p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
    p.font.name = TOKENS.TITLE_FONT
    p.font.bold = True

    if subtitle:
        sub_y = Inches(2.7 + h + 0.2)
        add_accent_bar(slide, Inches(4.8), sub_y, Inches(1.4), Inches(0.05), TOKENS.ACCENT)
        box = slide.shapes.add_textbox(Inches(4.8), sub_y + Inches(0.12), Inches(8.5), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = TOKENS.COVER_SUBTITLE
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.SECONDARY_GRAY))
        p.font.name = TOKENS.BODY_FONT

    add_footer(slide, slide_number, course_title)
    return slide


def build_content_slide(prs, topic, bullets, slide_number, course_title="", is_continued=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide)

    layout = Layout(slide, course_title, slide_number)
    display = f"{topic} (continued)" if is_continued else topic
    render_title(layout, display)
    render_accent_bar(layout, 1.0, TOKENS.ACCENT)
    render_bullet_card(layout, bullets)
    add_footer(slide, slide_number, course_title)
    return slide


def build_example_slide(prs, topic, example, slide_number, course_title=""):
    if not example:
        return None
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide)

    layout = Layout(slide, course_title, slide_number)
    render_title(layout, f"Example: {topic}")
    render_accent_bar(layout, 1.0, TOKENS.WARNING)

    paragraphs = ContentSplitter.to_paragraphs(example, max_chars=400)
    pad = to_inches(TOKENS.CARD_PADDING)
    text_w = layout.width_inches - 2 * pad - 0.15

    text_h = 0.0
    for para in paragraphs:
        text_h += TextMetrics.height(para, TOKENS.BODY_TEXT.pt, text_w, 1.35) + 0.12
    if paragraphs:
        text_h -= 0.12

    label_h = TextMetrics.line_h(11, 1.3)
    card_h = max(1.8, text_h + label_h + 2 * pad + 0.2)
    max_h = layout.remaining() - to_inches(TOKENS.GAP_MEDIUM)
    card_h = min(card_h, max_h)

    add_rounded_rect(slide, layout.left + Inches(0.12), Inches(layout.y), Inches(layout.width_inches - 0.24), Inches(card_h), TOKENS.CARD_BG, TOKENS.LIGHT_GRAY)

    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, layout.left + Inches(0.12), Inches(layout.y), Inches(0.07), Inches(card_h))
    set_fill(border, TOKENS.WARNING)
    no_outline(border)

    lbl = slide.shapes.add_textbox(layout.left + Inches(0.35), Inches(layout.y + 0.12), Inches(2), Inches(0.25))
    p = lbl.text_frame.paragraphs[0]
    p.text = "EXAMPLE"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WARNING))
    p.font.name = TOKENS.BODY_FONT
    p.font.bold = True

    txt_h = min(Inches(layout.remaining()), Inches(5.0))
    txt = slide.shapes.add_textbox(layout.left + Inches(0.35), Inches(layout.y + 0.38), Inches(layout.width_inches - 0.6), txt_h)
    tf = txt.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para
        p.font.size = TOKENS.BODY_TEXT
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
        p.font.name = TOKENS.BODY_FONT
        p.font.italic = True
        p.space_after = Pt(10)
        p.line_spacing = 1.35

    layout.advance(card_h + to_inches(TOKENS.GAP_MEDIUM))
    add_footer(slide, slide_number, course_title)
    return slide


def build_key_terms_slide(prs, key_terms, slide_number, course_title=""):
    if not key_terms:
        return None
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide)

    layout = Layout(slide, course_title, slide_number)
    render_title(layout, "Key Terms")
    render_accent_bar(layout, 1.0, TOKENS.ACCENT)

    card_w = Inches(5.6)
    gap_x = Inches(0.3)
    cols = 2
    pad = to_inches(TOKENS.CARD_PADDING)
    items = list(key_terms.items())

    for i, (term, definition) in enumerate(items):
        row = i // cols
        col = i % cols
        x = layout.left + col * (card_w + gap_x)

        def_w = to_inches(card_w) - 2 * pad - 1.9
        def_h = TextMetrics.height(definition, TOKENS.BODY_SMALL.pt, def_w, 1.35)
        term_h = TextMetrics.line_h(TOKENS.BODY_SMALL.pt, 1.3)
        card_h = max(1.0, def_h + term_h + 2 * pad + 0.1)
        y = layout.y + row * (card_h + 0.15)

        add_rounded_rect(slide, x, Inches(y), card_w, Inches(card_h), TOKENS.CARD_BG, TOKENS.LIGHT_GRAY)

        badge_w = min(Inches(1.8), Inches(len(term) * 0.09 + 0.3))
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.10), Inches(y + 0.15), badge_w, Inches(0.36))
        set_fill(badge, TOKENS.PRIMARY)
        no_outline(badge)

        tbox = slide.shapes.add_textbox(x + Inches(0.10), Inches(y + 0.17), badge_w, Inches(0.32))
        p = tbox.text_frame.paragraphs[0]
        p.text = term
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WHITE))
        p.font.name = TOKENS.BODY_FONT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        dbox = slide.shapes.add_textbox(x + badge_w + Inches(0.15), Inches(y + 0.12), card_w - badge_w - Inches(0.35), Inches(card_h - 0.24))
        tf = dbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = definition
        p.font.size = TOKENS.BODY_SMALL
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
        p.font.name = TOKENS.BODY_FONT
        p.line_spacing = 1.3

    rows = math.ceil(len(items) / cols)
    if rows > 0:
        layout.y += rows * 1.15
    add_footer(slide, slide_number, course_title)
    return slide


def build_summary_slide(prs, summary, slide_number, course_title=""):
    if not summary:
        return None
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide)

    layout = Layout(slide, course_title, slide_number)
    render_title(layout, "Lesson Summary")
    render_accent_bar(layout, 1.2, TOKENS.SUCCESS)

    paragraphs = ContentSplitter.to_paragraphs(summary, max_chars=400)
    pad = to_inches(TOKENS.CARD_PADDING)
    text_w = layout.width_inches - 2 * pad - 0.7

    text_h = 0.0
    for para in paragraphs:
        text_h += TextMetrics.height(para, TOKENS.BODY_TEXT.pt, text_w, 1.35) + 0.12
    if paragraphs:
        text_h -= 0.12

    card_h = max(1.8, text_h + 2 * pad + 0.2)
    max_h = layout.remaining() - to_inches(TOKENS.GAP_MEDIUM)
    card_h = min(card_h, max_h)

    add_rounded_rect(slide, layout.left, Inches(layout.y), Inches(layout.width_inches), Inches(card_h), TOKENS.CARD_BG, TOKENS.LIGHT_GRAY)

    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, layout.left + Inches(0.22), Inches(layout.y + 0.22), Inches(0.48), Inches(0.48))
    set_fill(icon, TOKENS.SUCCESS)
    no_outline(icon)

    itxt = slide.shapes.add_textbox(layout.left + Inches(0.22), Inches(layout.y + 0.25), Inches(0.48), Inches(0.42))
    p = itxt.text_frame.paragraphs[0]
    p.text = "✓"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WHITE))
    p.alignment = PP_ALIGN.CENTER

    txt_h = min(Inches(layout.remaining()), Inches(5.0))
    txt = slide.shapes.add_textbox(layout.left + Inches(0.85), Inches(layout.y + 0.15), Inches(layout.width_inches - 1.05), txt_h)
    tf = txt.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para
        p.font.size = TOKENS.BODY_TEXT
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT))
        p.font.name = TOKENS.BODY_FONT
        p.space_after = Pt(10)
        p.line_spacing = 1.35

    layout.advance(card_h + to_inches(TOKENS.GAP_MEDIUM))
    add_footer(slide, slide_number, course_title)
    return slide


def build_quiz_slide(prs, questions, slide_number, course_title=""):
    if not questions:
        return None
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide)

    layout = Layout(slide, course_title, slide_number)
    render_title(layout, "Check Your Understanding")
    render_accent_bar(layout, 1.2, TOKENS.ACCENT)
    render_numbered_card(layout, questions)
    add_footer(slide, slide_number, course_title)
    return slide


def build_references_slide(prs, papers, slide_number, course_title=""):
    if not papers:
        return None
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide)

    layout = Layout(slide, course_title, slide_number)
    render_title(layout, "References")
    render_accent_bar(layout, 1.0, TOKENS.ACCENT)

    for i, paper in enumerate(papers):
        if not isinstance(paper, dict):
            continue

        title = paper.get("title", "")
        authors = paper.get("authors", "")
        year = paper.get("year", "")
        lines = [title]
        if authors:
            lines.append(authors)
        if year:
            lines.append(f"({year})")

        pad = to_inches(TOKENS.CARD_PADDING)
        text_w = layout.width_inches - 2 * pad - 0.5
        text_h = 0.0
        for j, line in enumerate(lines):
            fs = TOKENS.BODY_TEXT.pt if j == 0 else TOKENS.BODY_SMALL.pt
            text_h += TextMetrics.height(line, fs, text_w, 1.3) + 0.06
        if lines:
            text_h -= 0.06

        card_h = max(0.8, text_h + 2 * pad + 0.1)
        max_h = layout.remaining() - to_inches(TOKENS.GAP_MEDIUM)

        if card_h > max_h and i > 0:
            add_footer(layout.slide, slide_number, course_title)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_header(slide)
            layout = Layout(slide, course_title, slide_number + 1)
            slide_number += 1
            render_title(layout, "References (continued)")
            render_accent_bar(layout, 1.0, TOKENS.ACCENT)
            card_h = min(card_h, layout.remaining() - to_inches(TOKENS.GAP_MEDIUM))

        add_rounded_rect(layout.slide, layout.left, Inches(layout.y), Inches(layout.width_inches), Inches(card_h), TOKENS.CARD_BG, TOKENS.LIGHT_GRAY)

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, layout.left + Inches(0.10), Inches(layout.y + 0.15), Inches(0.35), Inches(0.35))
        set_fill(badge, TOKENS.PRIMARY)
        no_outline(badge)

        num = slide.shapes.add_textbox(layout.left + Inches(0.10), Inches(layout.y + 0.17), Inches(0.35), Inches(0.31))
        p = num.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WHITE))
        p.font.name = TOKENS.BODY_FONT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        txt_h = min(Inches(layout.remaining()), Inches(3.0))
        txt = slide.shapes.add_textbox(layout.left + Inches(0.55), Inches(layout.y + 0.10), Inches(layout.width_inches - 0.7), txt_h)
        tf = txt.text_frame
        tf.word_wrap = True
        for j, line in enumerate(lines):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            if j == 0:
                p.font.size = TOKENS.BODY_TEXT
                p.font.bold = True
            else:
                p.font.size = TOKENS.BODY_SMALL
                p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.SECONDARY_GRAY))
            p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.DARK_TEXT if j == 0 else TOKENS.SECONDARY_GRAY))
            p.font.name = TOKENS.BODY_FONT
            p.line_spacing = 1.25

        layout.advance(card_h + to_inches(TOKENS.GAP_SMALL))

    add_footer(layout.slide, slide_number, course_title)
    return slide


def build_closing_slide(prs, course_title, slide_number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), TOKENS.SLIDE_WIDTH, TOKENS.SLIDE_HEIGHT)
    set_fill(bg, TOKENS.PRIMARY)
    no_outline(bg)

    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(1), Inches(3), Inches(3))
    set_fill(c1, TOKENS.LIGHTER_BLUE)
    no_outline(c1)

    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(4.5), Inches(2), Inches(2))
    set_fill(c2, TOKENS.VERY_LIGHT_BLUE)
    no_outline(c2)

    box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(56)
    p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WHITE))
    p.font.name = TOKENS.TITLE_FONT
    p.font.bold = True

    add_accent_bar(slide, Inches(1.5), Inches(3.5), Inches(1.4), Inches(0.05), TOKENS.ACCENT)

    qbox = slide.shapes.add_textbox(Inches(1.5), Inches(3.7), Inches(7), Inches(0.5))
    p = qbox.text_frame.paragraphs[0]
    p.text = "Questions?"
    p.font.size = TOKENS.COVER_SUBTITLE
    p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WHITE))
    p.font.name = TOKENS.BODY_FONT

    if course_title:
        cbox = slide.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(7), Inches(0.4))
        p = cbox.text_frame.paragraphs[0]
        p.text = course_title
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*hex_to_rgb(TOKENS.WHITE))
        p.font.name = TOKENS.BODY_FONT
        p.font.italic = True

    add_footer(slide, slide_number, course_title)
    return slide


# =============================================================================
# MAIN EXPORTER
# =============================================================================

def export_course_to_pptx_per_lesson(course: Dict, output_dir: str) -> List[str]:
    """Export a structured course to one PowerPoint file per lesson."""
    os.makedirs(output_dir, exist_ok=True)
    generated_files = []

    course_title = course.get("course_title", course.get("title", "Course"))
    learning_objectives = course.get("learning_objectives", [])
    modules = course.get("modules", [])

    for module_idx, module in enumerate(modules):
        if not isinstance(module, dict):
            continue

        module_title = module.get("module_title", f"Module {module_idx + 1}")
        difficulty = module.get("difficulty", "")
        problem = module.get("problem_addressed", "")
        solution = module.get("solution_approach", "")
        papers = module.get("based_on_papers", [])
        lessons = module.get("lessons", [])

        for lesson_idx, lesson in enumerate(lessons):
            if not isinstance(lesson, dict):
                continue

            lesson_title = lesson.get("lesson_title", f"Lesson {lesson_idx + 1}")
            sections = lesson.get("sections", [])
            summary = lesson.get("summary", "")
            check_understanding = lesson.get("check_understanding", [])

            prs = Presentation()
            prs.slide_width = TOKENS.SLIDE_WIDTH
            prs.slide_height = TOKENS.SLIDE_HEIGHT

            slide_number = 1

            build_title_slide(prs, course_title, lesson_title, module_title, slide_number)
            slide_number += 1

            if learning_objectives and lesson_idx == 0:
                build_objectives_slide(prs, learning_objectives, slide_number, course_title)
                slide_number += 1

            if lesson_idx == 0 and (problem or solution or difficulty):
                build_context_slide(prs, problem, solution, difficulty, slide_number, course_title)
                slide_number += 1

            for section in sections:
                if not isinstance(section, dict):
                    continue

                topic = section.get("topic", "")
                explanation = section.get("explanation", "")
                example = section.get("example_or_evidence", "")
                key_terms = section.get("key_terms", {})

                if not topic and not explanation:
                    continue

                if topic:
                    build_divider_slide(prs, topic, "", slide_number, course_title)
                    slide_number += 1

                if explanation:
                    groups = ContentSplitter.to_bullets(
                        explanation,
                        max_bullets=TOKENS.MAX_BULLETS_PER_SLIDE,
                        max_chars=TOKENS.MAX_CHARS_PER_BULLET
                    )
                    for i, bullets in enumerate(groups):
                        if not bullets:
                            continue
                        build_content_slide(prs, topic, bullets, slide_number, course_title, is_continued=(i > 0))
                        slide_number += 1

                if example:
                    build_example_slide(prs, topic, example, slide_number, course_title)
                    slide_number += 1

                if key_terms and isinstance(key_terms, dict):
                    chunks = ContentSplitter.chunk_dict(key_terms, per_slide=TOKENS.MAX_TERMS_PER_SLIDE)
                    for chunk in chunks:
                        build_key_terms_slide(prs, chunk, slide_number, course_title)
                        slide_number += 1

            if summary:
                build_summary_slide(prs, summary, slide_number, course_title)
                slide_number += 1

            if check_understanding:
                if isinstance(check_understanding, list):
                    questions = [str(q) for q in check_understanding if q]
                elif isinstance(check_understanding, str):
                    questions = [check_understanding]
                else:
                    questions = []

                if questions:
                    q_chunks = ContentSplitter.chunk_list(questions, per_slide=TOKENS.MAX_QUESTIONS_PER_SLIDE)
                    for chunk in q_chunks:
                        build_quiz_slide(prs, chunk, slide_number, course_title)
                        slide_number += 1

            if lesson_idx == len(lessons) - 1 and papers:
                build_references_slide(prs, papers, slide_number, course_title)
                slide_number += 1

            is_last = (module_idx == len(modules) - 1 and lesson_idx == len(lessons) - 1)
            if is_last:
                build_closing_slide(prs, course_title, slide_number)
                slide_number += 1

            safe_title = re.sub(r'[^\w\s-]', '', lesson_title).strip().replace(" ", "_")
            filename = f"{module_idx + 1:02d}_{lesson_idx + 1:02d}_{safe_title}.pptx"
            filepath = os.path.join(output_dir, filename)
            prs.save(filepath)
            generated_files.append(filepath)

    return generated_files


# Backward compatibility
export_lesson_to_pptx = export_course_to_pptx_per_lesson
